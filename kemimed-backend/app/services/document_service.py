import io
from loguru import logger


def parse_pdf(file_bytes: bytes) -> tuple[str, int]:
    """Extract text and page count from PDF bytes."""
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        pages = len(reader.pages)
        text = "\n".join(
            page.extract_text() or "" for page in reader.pages
        )
        return text.strip(), pages
    except Exception as e:
        logger.error(f"PDF parse error: {e}")
        return "", 0


def parse_pptx(file_bytes: bytes) -> tuple[str, int]:
    """Extract text from PowerPoint bytes."""
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts).strip(), len(prs.slides)
    except Exception as e:
        logger.error(f"PPTX parse error: {e}")
        return "", 0


def parse_docx(file_bytes: bytes) -> tuple[str, int]:
    """Extract text from DOCX bytes."""
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        text = "\n".join(p.text for p in doc.paragraphs)
        return text.strip(), len(doc.paragraphs)
    except Exception as e:
        logger.error(f"DOCX parse error: {e}")
        return "", 0
